import openai
import os
from io import BytesIO
from pdfminer.high_level import extract_text
from docx import Document
import pptx
from pptx import Presentation
import latexcodec
import requests
import shutil
import subprocess
import argparse
from PyPDF2 import PdfReader
from PIL import Image
import fitz
import pytesseract
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity


openai.api_key = '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'

def parse_pdf(pdf_file):
    return extract_text(pdf_file)

def parse_docx(docx_file):
    doc = Document(docx_file)
    text = ''
    for paragraph in doc.paragraphs:
        text += paragraph.text + '\n'
    return text

def parse_pptx(pptx_file):
    presentation = pptx.Presentation(pptx_file)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def parse_latex(latex_file):
    with open(latex_file, 'r', encoding='latex') as f:
        latex_content = f.read()
    return latex_content

def parse_document(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == '.pdf':
        return parse_pdf(file_path)
    elif file_extension == '.docx':
        return parse_docx(file_path)
    elif file_extension == '.pptx':
        return parse_pptx(file_path)
    elif file_extension == '.tex':
        return parse_latex(file_path)
    else:
        raise ValueError("Unsupported file type")

def ask_question(document_text, question):
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": f"Text:\n{document_text}\n\n"},
            {"role": "user", "content": f"{question}\n"}
        ],
        temperature=0.5,
        max_tokens=1000
    )
    return response.choices[0].message['content'].strip()


def extract_text_from_image(image):
    return pytesseract.image_to_string(image)

def extract_text_from_pdf_images(pdf_file):
    text = ''
    with fitz.open(pdf_file) as pdf:
        for page_number in range(len(pdf)):
            page = pdf.load_page(page_number)
            images = page.get_images(full=True)
            for img_info in images:
                base_image = pdf.extract_image(img_info[0])
                image_bytes = base_image["image"]
                image = Image.open(BytesIO(image_bytes))
                text += extract_text_from_image(image) + '\n'
    return text

def extract_text_from_docx_images(docx_file):
    text = ''
    doc = Document(docx_file)
    for rel in doc.part.rels:
        if "image" in doc.part.rels[rel].target_ref:
            image_part = doc.part.rels[rel].target_part
            image = Image.open(BytesIO(image_part.blob))
            text += extract_text_from_image(image) + '\n'
    return text

def extract_text_from_pptx_images(pptx_file):
    text = ''
    presentation = Presentation(pptx_file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                image = shape.image.blob
                image = Image.open(BytesIO(image))
                text += extract_text_from_image(image) + '\n'
    return text

def extract_text_from_images(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == '.pdf':
        return extract_text_from_pdf_images(file_path)
    elif file_extension == '.docx':
        return extract_text_from_docx_images(file_path)
    elif file_extension == '.pptx':
        return extract_text_from_pptx_images(file_path)
    else:
        raise ValueError("Unsupported file type")

def compute_similarity(query, document_text):
    vectorizer = TfidfVectorizer()
    vectors = vectorizer.fit_transform([query] + document_text.split('\n'))
    similarity_scores = cosine_similarity(vectors)
    return similarity_scores

def find_relevant_lines(document_text, answer, threshold=0.25):
    lines = document_text.split('\n')
    similarity_scores = compute_similarity(answer, document_text)

    relevant_lines = []
    for i, score in enumerate(similarity_scores[0][1:]):
        if score > threshold:
            relevant_lines.append({"line_number": i + 1, "similarity_score": score, "line_text": lines[i]})
    return relevant_lines

def main(file_path, question):
    file_extension = os.path.splitext(file_path)[1].lower()
    document_text = parse_document(file_path)
    if file_extension == '.tex':
        image_data = ""
    else:
        image_data = extract_text_from_images(file_path)
    combined_data = {
        "document_text": document_text,
        "image_data": image_data
    }
    answer = ask_question(combined_data, question)
    relevant_lines = find_relevant_lines(document_text + image_data, answer)

    print("Answer:")
    print("-" * 50)
    print(answer)
    print("\n\nRelevant Lines:")
    print("-" * 50)
    for line_info in relevant_lines:
        print(f"\n👉 Line {line_info['line_number']}: Similarity Score - {line_info['similarity_score']:.4f}")
        print(f"    {line_info['line_text']}")
        print("-" * 50)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Run python script.")
    parser.add_argument('file', type=str, help='File Location')
    parser.add_argument('question', type=str, help='Input Question')
    args = parser.parse_args()
    main(args.file, args.question)